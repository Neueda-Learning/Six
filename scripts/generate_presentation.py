"""生成课程项目答辩用 PowerPoint（Payments Processing System）。
运行一次即可在仓库根目录生成 Payments-Processing-Presentation.pptx，可按需重新运行覆盖生成。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY = RGBColor(0x25, 0x63, 0xEB)   # 品牌蓝，与前端 UI 主色一致
DARK = RGBColor(0x1F, 0x23, 0x29)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_title(slide, text, size=40, color=DARK, top=Inches(0.5), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(0.7), top, Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return box


def add_subtitle(slide, text, top=Inches(1.6), size=20, color=GRAY):
    box = slide.shapes.add_textbox(Inches(0.7), top, Inches(12), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, top=Inches(1.8), left=Inches(0.9), width=Inches(11.5),
                 height=Inches(5.2), size=20, color=DARK, bullet_char="• "):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        run = p.add_run()
        run.text = ("    " * level) + (bullet_char if level == 0 else "- ") + text
        run.font.size = Pt(size - level * 2)
        run.font.color.rgb = color
        run.font.bold = (level == 0)
    return box


def section_slide(title, subtitle=None):
    slide = add_slide()
    set_background(slide, PRIMARY)
    add_title(slide, title, size=44, color=WHITE, top=Inches(3.0), align=PP_ALIGN.LEFT)
    if subtitle:
        add_subtitle(slide, subtitle, top=Inches(4.0), color=RGBColor(0xDB, 0xEA, 0xFE))
    return slide


def content_slide(title, bullets, subtitle=None):
    slide = add_slide()
    set_background(slide, WHITE)
    add_title(slide, title, size=32)
    top = Inches(1.5)
    if subtitle:
        add_subtitle(slide, subtitle, top=Inches(1.3))
        top = Inches(2.1)
    add_bullets(slide, bullets, top=top)
    return slide


# 1. 封面
slide = add_slide()
set_background(slide, PRIMARY)
add_title(slide, "Payments Processing System", size=48, color=WHITE, top=Inches(2.6))
add_subtitle(slide, "支付处理生命周期管理系统 —— 课程项目答辩", top=Inches(3.6), color=RGBColor(0xDB, 0xEA, 0xFE))
add_subtitle(slide, "Spring Boot 3 + MyBatis-Plus + MySQL  |  Vue 3 + Element Plus + vue-i18n",
             top=Inches(4.3), size=16, color=RGBColor(0xDB, 0xEA, 0xFE))

# 2. Agenda
content_slide("目录 Agenda", [
    "项目背景与目标",
    "支付生命周期与状态机设计",
    "系统架构与技术栈",
    "核心功能演示（Live Demo）",
    "我的贡献：五大功能模块详解",
    "遇到的挑战与解决方案",
    "测试覆盖与质量保障",
    "未来展望",
])

# 3. 项目背景
content_slide("项目背景与目标", [
    "目标：设计并实现一个处理金融支付完整生命周期的系统",
    "核心要求：创建、验证、处理、完成/失败的全流程管理",
    "维护完整的状态变更 Audit Trail（审计时间线）",
    "无需认证、单用户假设，不对接真实支付网关（内部模拟）",
    "技术栈由团队自主选择：Spring Boot + Vue",
])

# 4. 状态机设计
content_slide("支付生命周期与状态机设计", [
    "五种状态：CREATED → VALIDATED → SENT → COMPLETED",
    "任意阶段均可流转至 FAILED（终态）",
    "合法流转由 PaymentStateMachine 集中维护，杜绝跳级/逆向",
    ("CREATED → VALIDATED：三层校验（金额 / 币种 / 账户）+ 余额充足性校验", 1),
    ("VALIDATED → SENT，SENT → COMPLETED：正常推进", 1),
    ("任意阶段 → FAILED：携带 errorCode + errorMessage", 1),
])

# 5. 架构
content_slide("系统架构与技术栈", [
    "前端：Vue 3 + Element Plus + vue-router + vue-i18n（三语言：中/英/德）",
    "后端：Spring Boot 3 + MyBatis-Plus + MySQL，分层架构 Controller/Service/Validator/StateMachine",
    "统一响应体：{ success, data, errorCode, message }",
    "全局异常处理：@RestControllerAdvice 统一转换错误响应",
    "接口文档：Swagger / OpenAPI",
    "定时任务：PaymentAutoTransitionScheduler 模拟真实网络异步处理延迟",
])

# 6. 功能模块概览
content_slide("核心功能模块概览", [
    "支付创建（幂等键防重复提交）",
    "支付详情查看 + 失败错误详情展示",
    "支付状态历史时间线（Audit Trail）",
    "支付列表检索、按状态筛选、分页",
    "回收站（软删除 / 恢复 / 永久删除，30 天保留期）",
    "手动状态流转接口（课程演示 / 测试用）",
])

# 7. 我的贡献 分割页
section_slide("我的贡献", "五个由我独立设计并实现的功能模块")

# 8. 贡献1 i18n
content_slide("贡献 1：多语言国际化（i18n）", [
    "引入 vue-i18n，实现中文 / English / Deutsch 三语言切换",
    "首次访问按浏览器语言自动探测，其后持久化到 localStorage",
    "el-config-provider 联动 Element Plus 内置组件语言（分页器等）",
    "金额 / 日期按语言环境本地化格式化（n() / d()）",
    "错误码前端本地化字典，优先于后端固定语言 message 展示",
    "仅改动前端代码，未触碰后端接口",
])

# 9. 贡献2 auto transition
content_slide("贡献 2：自动状态推进 + 随机失败模拟", [
    "问题：原自动推进逻辑必然走向 COMPLETED，无法演示 FAILED 分支",
    "方案：新增可配置的 failure-probability（默认 20%）",
    "按当前所处阶段选用贴近场景的错误码（PROCESSING_ERROR / NETWORK_ERROR）",
    "单次 20% 概率下，全流程整体失败率约 49%，同时演示两种真实分支",
    "与余额校验协同：真实业务规则优先于随机演示逻辑",
])

# 10. 贡献3 余额校验
content_slide("贡献 3：账户余额充足性校验", [
    "需求：只做余额充足性判断，不扣款、不涉及并发扣减，余额自始至终不变",
    "新增 balance 字段（独立迁移脚本，不改动 schema.sql / data.sql）",
    "hasSufficientBalance()：只读查询，无任何 UPDATE 语句",
    "只在 CREATED → VALIDATED 这一步触发校验",
    "余额不足时，实际执行的目标状态被强制改写为 FAILED / INSUFFICIENT_FUNDS",
    "手动接口与自动调度两处入口统一复用同一段判断逻辑",
])

# 11. 贡献4 轮询
content_slide("贡献 4：列表页 / 详情页自动轮询", [
    "问题：后端每 5 秒自动推进状态，前端不刷新则用户看不到变化",
    "详情页：单笔支付非终态时，每 5 秒静默刷新，到终态自动停止",
    "列表页（新增）：只要当前页存在任意非终态记录就继续轮询",
    "静默刷新不显示整页 loading 遮罩，避免闪烁",
    "onUnmounted 清理定时器，离开页面不产生后台请求",
])

# 12. 贡献5 测试
content_slide("贡献 5：测试用例文档 + JUnit 5 单元测试", [
    "test-cases.md：14 个章节，96 个测试用例，覆盖全部核心场景",
    "主动标注 5 处代码审查发现的“待确认事项”，不回避已知缺口",
    "PaymentValidatorTest：金额 / 账户 / 余额校验，15 个用例",
    "PaymentStateMachineTest：参数化测试一次覆盖 25 种状态流转组合",
    "PaymentServiceImplTest：状态解析 + 乐观锁冲突场景（如实记录代码现状）",
    "全部使用 Mockito，隔离数据库与后台调度任务，测试稳定可重复",
])

# 13. Demo
content_slide("Live Demo 演示提纲", [
    "1. 创建一笔支付（展示幂等键、金额/币种/账户校验）",
    "2. 切换语言（中/英/德），展示界面文案与金额/日期格式联动变化",
    "3. 打开详情页，展示自动轮询：状态从 CREATED 自动推进",
    "4. 演示一笔余额不足的支付：强制转为 FAILED / INSUFFICIENT_FUNDS",
    "5. 查看状态历史时间线（Audit Trail）",
    "6. 演示回收站：软删除 → 恢复 / 永久删除",
])

# 14. 挑战
content_slide("遇到的挑战与解决方案", [
    "挑战：后台定时任务与集成测试断言冲突（flaky test）",
    ("解决：改用 Mockito 纯单元测试，不依赖真实 Spring 容器与数据库", 1),
    "挑战：多语言切换后，Element Plus 内置组件文案未同步",
    ("解决：引入 el-config-provider 统一注入语言配置", 1),
    "挑战：随机失败模拟与真实余额校验的错误码语义冲突",
    ("解决：明确优先级——真实业务规则先于随机演示逻辑判断", 1),
])

# 15. 未来展望
content_slide("未来展望", [
    "补充 Controller 层 @WebMvcTest，提升整体测试覆盖率至 80%+",
    "引入可 mock 的 Clock，解决回收站保留期边界测试的时间注入问题",
    "为乐观锁冲突场景补充真实的异常检测与友好错误码",
    "多语言支持更多语种，補充德语专业术语校对",
    "接入真实 CI/CD，将 JaCoCo 覆盖率检查纳入流水线门禁",
])

# 16. 致谢
slide = add_slide()
set_background(slide, PRIMARY)
add_title(slide, "谢谢聆听", size=44, color=WHITE, top=Inches(3.0))
add_subtitle(slide, "Q & A", top=Inches(4.0), size=24, color=RGBColor(0xDB, 0xEA, 0xFE))

prs.save("Payments-Processing-Presentation.pptx")
print("Saved Payments-Processing-Presentation.pptx with", len(prs.slides), "slides")
